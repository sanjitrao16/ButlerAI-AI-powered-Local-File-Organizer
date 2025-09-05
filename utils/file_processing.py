import fitz #PyMuPDF
import docx
import pandas as pd
from pptx import Presentation

def process_txt_file(file_path):
    max_chars_to_read = 2500
    try:
        with open(file_path,"r",encoding="utf-8",errors="ignore") as f:
            text = f.read(max_chars_to_read)
            return text
    except Exception as e:
        print(f"[Butler AI] Error reading TXT/MD file {file_path}: {e}")
        return None

def process_doc_file(file_path):
    try:
        word_doc = docx.Document(file_path)
        full_text = [para.text for para in word_doc.paragraphs]
        return "\n".join(full_text)
    except Exception as e:
        print(f"[Butler AI] Error reading word document file {file_path}: {e}")
        return None

def process_ppt_file(file_path):
    try:
        ppt_file = Presentation(file_path)
        full_text = []
        for slide in ppt_file.slides:
            for shape in slide.shapes:
                if hasattr(shape,"text"):
                    full_text.append(shape.text)
        return "\n".join(full_text)
    except Exception as e:
        print(f"[Butler AI] Error reading presentation file {file_path}: {e}")
        return None

def process_pdf_file(file_path):
    try:
        pdf_file = fitz.open(file_path)
        max_pages_to_read = 4
        full_text = []
        for page_num in range(min(max_pages_to_read, len(pdf_file))):
            page = pdf_file.load_page(page_num)
            full_text.append(page.get_text())
        pdf_content = "\n".join(full_text)
        return pdf_content
    except Exception as e:
        print(f"[Butler AI] Error reading PDF file {file_path}: {e}")
        return None

def process_spreadsheet_file(file_path):
    try:
        if file_path.lower().endswith(".csv"):
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path)
        text = df.to_string()
        return text
    except Exception as e:
        print(f"[Butler AI] Error reading spreadsheet file {file_path}: {e}")
        return None
