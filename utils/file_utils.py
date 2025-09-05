import os
import shutil
import sys
import re
import concurrent.futures
import fitz #PyMuPDF
import docx
import pandas as pd
from pptx import Presentation

def process_txt_file(file_path):
    max_chars_to_read = 2500
    try:
        with open(file_path,"r",encoding="utf-8",errors="ignore") as f:
            text = f.read(max_chars_to_read)
            return text
    except Exception as e:
        print(f"[Butler AI] Error reading TXT/MD file {file_path}: {e}")
        return None

def process_doc_file(file_path):
    try:
        word_doc = docx.Document(file_path)
        full_text = [para.text for para in word_doc.paragraphs]
        return "\n".join(full_text)
    except Exception as e:
        print(f"[Butler AI] Error reading word document file {file_path}: {e}")
        return None

def process_ppt_file(file_path):
    try:
        ppt_file = Presentation(file_path)
        full_text = []
        for slide in ppt_file.slides:
            for shape in slide.shapes:
                if hasattr(shape,"text"):
                    full_text.append(shape.text)
        return "\n".join(full_text)
    except Exception as e:
        print(f"[Butler AI] Error reading presentation file {file_path}: {e}")
        return None

def process_pdf_file(file_path):
    try:
        pdf_file = fitz.open(file_path)
        max_pages_to_read = 4
        full_text = []
        for page_num in range(min(max_pages_to_read, len(pdf_file))):
            page = pdf_file.load_page(page_num)
            full_text.append(page.get_text())
        pdf_content = "\n".join(full_text)
        return pdf_content
    except Exception as e:
        print(f"[Butler AI] Error reading PDF file {file_path}: {e}")
        return None

def process_spreadsheet_file(file_path):
    try:
        if file_path.lower().endswith(".csv"):
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path)
        text = df.to_string()
        return text
    except Exception as e:
        print(f"[Butler AI] Error reading spreadsheet file {file_path}: {e}")
        return None


def extract_text_content(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext in [".txt",".md"]:
            return (file_path, process_txt_file(file_path))
        elif ext in [".doc",".docx"]:
            return (file_path, process_doc_file(file_path))
        elif ext in [".ppt",".pptx"]:
            return (file_path, process_ppt_file(file_path))
        elif ext in [".pdf"]:
            return (file_path, process_pdf_file(file_path))
        elif ext in [".csv",".xls",".xlsx"]:
            return (file_path, process_spreadsheet_file(file_path))
    except Exception as e:
        print(f"[Butler AI] File type not supported {file_path}: {e}")


def display_directory_tree(dir_path, prefix=""):
    contents = sorted(os.listdir(dir_path))
    pointers = ['├── '] * (len(contents) - 1) + ['└── ']
    for pointer, name in zip(pointers, contents):
        path = os.path.join(dir_path, name)
        print(prefix + pointer + name)
        if os.path.isdir(path):
            extension = '│   ' if pointer == '├── ' else '    '
            display_directory_tree(path, prefix + extension)


def separate_files_by_type(dir_path):
    '''
    Text-File formats: .pdf, .docx, .pptx, .txt, .csv, .md, .xlsx
    Image-File formats: .jpg, .jpeg, .png, .webp, .gif
    Video-File formats: .mp4
    Audio-File formats: .mp3
    
    '''
    text_exts = (".csv",".doc",".docx",".md",".pdf",".ppt",".pptx",".txt",".xls",".xlsx")
    image_exts = (".gif",".jpeg",".jpg",".png",".webp")
    video_exts = (".mp4")
    audio_exts = (".mp3")

    files = [
        os.path.join(dir_path,name) for name in os.listdir(dir_path) if os.path.isfile(os.path.join(dir_path,name))
    ]

    text_files = [file for file in files if os.path.splitext(file)[1].lower() in text_exts]
    image_files = [file for file in files if os.path.splitext(file)[1].lower() in image_exts]
    video_files = [file for file in files if os.path.splitext(file)[1].lower() in video_exts]
    audio_files = [file for file in files if os.path.splitext(file)[1].lower() in audio_exts]
    
    return video_files,audio_files,image_files,text_files

def process_text_files(text_files,workers=4):
    results = {}
    with concurrent.futures.ProcessPoolExecutor(max_workers=workers) as executor:
        for file_path,content in executor.map(extract_text_content,text_files):
            results[file_path] = content
    return results